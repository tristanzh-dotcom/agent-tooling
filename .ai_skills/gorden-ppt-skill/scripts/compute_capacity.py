#!/usr/bin/env python3
"""Compute geometry-accurate text capacity for every editable slot in a
template's detail.json, using the template.pptx for box size + font size.

Adds to each slot:
  - box_cm: [width_cm, height_cm]
  - font_size_pt: resolved effective size (no None — inheritance resolved)
  - chars_per_line: max CJK-equivalent units per line (visual-width units)
  - max_lines: max lines that fit in the box height
  - max_chars: capacity in CJK-equivalent units (chars_per_line * max_lines)
  - wrap: whether the box wraps text
  - autofit: whether the box auto-shrinks text on overflow

Adds to the template root:
  - type_scale: [{level, size_pt, count}] ranked largest->smallest
  and tags each slot with `level` (1 = largest size tier).

Capacity is measured in "visual width units" (vw): a CJK glyph = 1.0,
latin/digit ≈ 0.5, space ≈ 0.35. 1 vw unit ≈ 1 em ≈ font_size_pt wide.

Usage:
  python3 compute_capacity.py <detail.json> <template.pptx> [--in-place|-o out.json]
"""
from __future__ import annotations
import argparse
import json
import math
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE

EMU_PER_CM = 360000
PT_PER_CM = 28.3465
A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


# ---------------- text visual width ----------------
def vw_of(s: str) -> float:
    w = 0.0
    for c in s:
        if "\u4e00" <= c <= "\u9fff" or "\u3000" <= c <= "\u303f" or "\uff00" <= c <= "\uffef":
            w += 1.0          # CJK / fullwidth
        elif c == " ":
            w += 0.35
        elif c.isascii():
            w += 0.5          # latin / digit / punctuation
        else:
            w += 0.8
    return w


# ---------------- shape lookup (recurse groups) ----------------
def find_shape(shapes, shape_id: int):
    for shape in shapes:
        if shape.shape_id == shape_id:
            return shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            f = find_shape(shape.shapes, shape_id)
            if f is not None:
                return f
    return None


# ---------------- font size resolution ----------------
def _defRPr_sz_from_lststyle(el, level: int):
    """Look for a:lstStyle/a:lvlNpPr/a:defRPr@sz inside an element."""
    if el is None:
        return None
    lst = el.find(f"{A_NS}lstStyle") if not el.tag.endswith("lstStyle") else el
    # txBody path: el may be a txBody; search its lstStyle
    if lst is None:
        # try shape's txBody
        txBody = el.find(f".//{A_NS}txBody")
        if txBody is not None:
            lst = txBody.find(f"{A_NS}lstStyle")
    if lst is None:
        return None
    lvl_tag = f"{A_NS}lvl{level+1}pPr"
    lvl = lst.find(lvl_tag)
    if lvl is None:
        return None
    defRPr = lvl.find(f"{A_NS}defRPr")
    if defRPr is not None and defRPr.get("sz"):
        return int(defRPr.get("sz")) / 100.0
    return None


def _master_txstyle_sz(master, ph_type, level: int):
    """master txStyles: title/body/other style defRPr sz by level."""
    txStyles = master._element.find(f"{A_NS}txStyles")
    if txStyles is None:
        return None
    # choose style bucket
    t = str(ph_type) if ph_type is not None else ""
    if "TITLE" in t or "CENTER_TITLE" in t or "ctrTitle" in t.lower():
        bucket = "titleStyle"
    elif "BODY" in t or "SUBTITLE" in t or "OBJECT" in t:
        bucket = "bodyStyle"
    else:
        bucket = "otherStyle"
    style = txStyles.find(f"{A_NS}{bucket}")
    if style is None:
        return None
    lvl = style.find(f"{A_NS}lvl{level+1}pPr")
    if lvl is None:
        return None
    defRPr = lvl.find(f"{A_NS}defRPr")
    if defRPr is not None and defRPr.get("sz"):
        return int(defRPr.get("sz")) / 100.0
    return None


def resolve_size_pt(shape, paragraph, slide, role: str) -> float:
    # 1) explicit run size already handled by caller; here resolve inherited
    # 2) paragraph-level font size
    try:
        if paragraph.font.size is not None:
            return paragraph.font.size.pt
    except Exception:
        pass
    level = getattr(paragraph, "level", 0) or 0

    # 3) placeholder inheritance: layout ph -> master ph -> master txStyles
    ph_type = None
    ph_idx = None
    try:
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            ph_idx = shape.placeholder_format.idx
    except Exception:
        pass

    if ph_idx is not None:
        layout = slide.slide_layout
        master = layout.slide_master
        for container in (layout, master):
            try:
                for ph in container.placeholders:
                    if ph.placeholder_format.idx == ph_idx:
                        sz = _defRPr_sz_from_lststyle(ph._element, level)
                        if sz:
                            return sz
            except Exception:
                pass
        sz = _master_txstyle_sz(master, ph_type, level)
        if sz:
            return sz

    # 4) role-based default
    r = role or ""
    if "主标题" in r or "页面标题" in r:
        return 32.0
    if "段落标题" in r or "副标题" in r or "小标题" in r:
        return 18.0
    if "正文" in r:
        return 13.0
    # 5) box-height heuristic (single line)
    try:
        h_cm = shape.height / EMU_PER_CM
        return max(10.0, min(40.0, h_cm * PT_PER_CM / 1.4))
    except Exception:
        return 14.0


# ---------------- capacity ----------------
# Calibrated against original (designer-fit) text across multiple templates:
# these constants keep the false-positive rate ~0-5% while still catching
# genuine 1.2x+ overflows.
H_MARGIN_CM = 0.25     # left+right text inset (combined)
V_MARGIN_CM = 0.08     # top+bottom text inset (combined)
LINE_HEIGHT = 1.0      # CJK body is typically single-spaced
TOLERANCE = 1.2        # model slack: allow 20% over raw geometry before flagging


def capacity_for(width_cm: float, height_cm: float, size_pt: float,
                 wrap: bool) -> tuple[int, int, int]:
    usable_w_pt = max(0.0, (width_cm - H_MARGIN_CM)) * PT_PER_CM
    usable_h_pt = max(0.0, (height_cm - V_MARGIN_CM)) * PT_PER_CM
    if size_pt <= 0:
        size_pt = 14.0
    cpl = max(1, math.floor(usable_w_pt / size_pt))       # vw units per line
    if not wrap:
        max_lines = 1
    else:
        max_lines = max(1, math.floor(usable_h_pt / (size_pt * LINE_HEIGHT)))
    cap = max(1, math.floor(cpl * max_lines * TOLERANCE))  # budget with slack
    return cpl, max_lines, cap


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("detail_json", type=Path)
    ap.add_argument("pptx", type=Path)
    ap.add_argument("-o", "--output", type=Path)
    args = ap.parse_args()

    detail = json.loads(args.detail_json.read_text(encoding="utf-8"))
    prs = Presentation(args.pptx)
    slides = list(prs.slides)

    sizes_seen: dict[float, int] = {}

    for page in detail["pages"]:
        sl_idx = page["slide_number"] - 1
        if sl_idx >= len(slides):
            continue
        slide = slides[sl_idx]
        for slot in page["text_slots"]:
            addr = slot["address"]
            shape = find_shape(slide.shapes, addr["shape_id"])
            if shape is None or not shape.has_text_frame:
                continue
            tf = shape.text_frame
            # geometry
            try:
                w_cm = shape.width / EMU_PER_CM
                h_cm = shape.height / EMU_PER_CM
            except Exception:
                continue
            slot["box_cm"] = [round(w_cm, 2), round(h_cm, 2)]

            # Degenerate box (e.g. auto-grow text box or unmeasurable group
            # child): geometry unreliable -> don't write capacity so the build
            # lint skips it; keep the existing heuristic max_chars as guidance.
            if w_cm < 1.0 or h_cm < 0.3:
                slot["capacity_unknown"] = True
                for k in ("chars_per_line", "max_lines"):
                    slot.pop(k, None)
                continue
            slot.pop("capacity_unknown", None)

            # font size: explicit on the addressed run/para, else resolve
            size_pt = slot.get("font_size_pt")
            pi = addr.get("paragraph", 0)
            para = tf.paragraphs[pi] if pi < len(tf.paragraphs) else tf.paragraphs[0]
            if not size_pt:
                ri = addr.get("run")
                run_sz = None
                if ri is not None and ri < len(para.runs):
                    rs = para.runs[ri].font.size
                    run_sz = rs.pt if rs is not None else None
                size_pt = run_sz or resolve_size_pt(shape, para, slide, slot.get("role", ""))
            # round to nearest 0.5pt: real design sizes are integral; fractional
            # values are heuristic artifacts that should collapse into one tier
            size_pt = round(size_pt * 2) / 2
            slot["font_size_pt"] = size_pt

            # wrap + autofit
            wrap = tf.word_wrap
            wrap = True if wrap is None else bool(wrap)
            slot["wrap"] = wrap
            autofit = False
            try:
                autofit = (tf.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE)
            except Exception:
                pass
            slot["autofit"] = autofit

            cpl, max_lines, cap = capacity_for(w_cm, h_cm, size_pt, wrap)
            slot["chars_per_line"] = cpl
            slot["max_lines"] = max_lines
            slot["max_chars"] = cap

            if slot.get("editable"):
                sizes_seen[slot["font_size_pt"]] = sizes_seen.get(slot["font_size_pt"], 0) + 1

    # type_scale: rank distinct sizes desc
    ranked = sorted(sizes_seen.items(), key=lambda kv: -kv[0])
    size_to_level = {sz: i + 1 for i, (sz, _) in enumerate(ranked)}
    detail["type_scale"] = [
        {"level": i + 1, "size_pt": sz, "slot_count": cnt}
        for i, (sz, cnt) in enumerate(ranked)
    ]
    for page in detail["pages"]:
        for slot in page["text_slots"]:
            sp = slot.get("font_size_pt")
            if sp in size_to_level:
                slot["level"] = size_to_level[sp]

    out = args.output or args.detail_json
    out.write_text(json.dumps(detail, ensure_ascii=False, indent=2), encoding="utf-8")
    n_slots = sum(len(p["text_slots"]) for p in detail["pages"])
    print(f"Wrote {out}  | {len(detail['pages'])} pages, {n_slots} slots, "
          f"{len(ranked)} size tiers")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
