#!/usr/bin/env python3
"""Build a new .pptx from a template by selecting slides and replacing text.

The script preserves all layout, fonts, colors, images and animations of the
source template. It only:
  1. Keeps the slides listed in ``selected_slides`` (in that order),
  2. Replaces the text of the runs identified by ``edits``.

USAGE
    python3 build_pptx.py <template.pptx> <edits.json> <output.pptx>
                         [--detail detail.json]   # for slot_id resolution
                         [--strict]               # fail if expected_text mismatch
                         [--dry-run]              # show planned edits, do nothing

EDITS JSON SCHEMA
{
  "template_slug": "simple-business",                 # optional, informational
  "selected_slides": [1, 2, 4, 7, 14, 16],            # 1-indexed in template
  "edits": [
    {
      "slide": 1,                                      # 1-indexed in template
      "slot_id": "title_en",                           # resolved via detail.json
      "new_text": "Annual Review"
    },
    {
      "slide": 1,
      "address": {"shape_id": 46, "paragraph": 1, "run": 0},  # explicit
      "expected_text": "工作计划模板",
      "new_text": "2026 年度工作计划"
    }
  ]
}

ADDRESS FORMS
    {"shape_id": <int>, "paragraph": <int>, "run": <int>}   # exact run
    {"shape_id": <int>, "paragraph": <int>}                  # whole paragraph
                                                              # (run 0 keeps fmt,
                                                              # others cleared)
"""
from __future__ import annotations
import argparse
import copy
import json
import math
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


SLIDE_R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"


# ------------------------------------------------------------------
# Slot resolution via detail.json
# ------------------------------------------------------------------
def load_slot_index(detail_path: Path | None) -> dict:
    """Return {(slide_number, slot_id): slot_dict}.

    slot_dict includes address + capacity metadata (max_chars, chars_per_line,
    max_lines, wrap, autofit, font_size_pt, level, role) when present so the
    overflow lint can use them.
    """
    if detail_path is None or not detail_path.exists():
        return {}
    detail = json.loads(detail_path.read_text(encoding="utf-8"))
    index: dict = {}
    for page in detail.get("pages", []):
        slide_number = page["slide_number"]
        for slot in page.get("text_slots", []):
            slot.setdefault("expected_text", slot.get("current_text"))
            index[(slide_number, slot["slot_id"])] = slot
    return index


# ------------------------------------------------------------------
# Overflow lint (text-fits-the-box check)
# ------------------------------------------------------------------
def _visual_width(s: str) -> float:
    w = 0.0
    for c in s:
        if "\u4e00" <= c <= "\u9fff" or "\u3000" <= c <= "\u303f" or "\uff00" <= c <= "\uffef":
            w += 1.0
        elif c == " ":
            w += 0.35
        elif c.isascii():
            w += 0.5
        else:
            w += 0.8
    return w


def check_overflow(new_text: str, meta: dict) -> tuple[bool, str]:
    """Return (fits, message). Uses chars_per_line / max_lines from detail.json.

    'fits' is True when the text is expected to stay within the box. Boxes with
    autofit (PowerPoint shrink-to-fit) are always treated as fitting (soft).
    """
    if meta.get("capacity_unknown"):
        return True, ""  # geometry unreliable for this box
    cpl = meta.get("chars_per_line")
    max_lines = meta.get("max_lines")
    cap = meta.get("max_chars")   # already includes the calibration tolerance
    if not cpl or not max_lines or not cap:
        return True, ""  # no capacity data → can't lint
    total_vw = _visual_width(new_text.replace("\n", ""))
    # primary test: visual width vs the (tolerance-adjusted) capacity budget
    if total_vw <= cap:
        return True, ""
    # estimate needed lines for a helpful message
    needed = sum(max(1, math.ceil(_visual_width(seg) / cpl))
                 for seg in (new_text.split("\n") or [new_text]))
    sz = meta.get("font_size_pt")
    msg = (f"视觉宽度 {total_vw:.0f} > 容量 {cap} "
           f"(约需 {needed} 行 / 可容 {max_lines} 行, {sz}pt, 每行≈{cpl}字)")
    if meta.get("autofit"):
        return True, "AUTOFIT " + msg + " — PowerPoint 会自动缩字"
    return False, msg


# ------------------------------------------------------------------
# Shape lookup (recurses into groups)
# ------------------------------------------------------------------
def find_shape(shapes, shape_id: int):
    for shape in shapes:
        if shape.shape_id == shape_id:
            return shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            found = find_shape(shape.shapes, shape_id)
            if found is not None:
                return found
    return None


# ------------------------------------------------------------------
# Text mutation that preserves run-level formatting
# ------------------------------------------------------------------
def apply_edit_to_shape(shape, address: dict, new_text: str, expected: str | None,
                       strict: bool) -> tuple[bool, str]:
    if not shape.has_text_frame:
        return False, f"shape {shape.shape_id} has no text frame"
    tf = shape.text_frame
    pi = address.get("paragraph", 0)
    if pi >= len(tf.paragraphs):
        return False, f"paragraph {pi} out of range (have {len(tf.paragraphs)})"
    para = tf.paragraphs[pi]
    runs = list(para.runs)
    if not runs:
        return False, f"paragraph {pi} has no runs"

    ri = address.get("run")
    if ri is None:
        current = "".join(r.text for r in runs)
        if expected is not None and current != expected and strict:
            return False, (
                f"expected_text mismatch at paragraph level: "
                f"have {current!r}, expected {expected!r}"
            )
        runs[0].text = new_text
        for r in runs[1:]:
            r.text = ""
        return True, f"replaced paragraph (kept run 0 formatting): {current!r} -> {new_text!r}"

    if ri >= len(runs):
        return False, f"run {ri} out of range (have {len(runs)})"
    current = runs[ri].text
    if expected is not None and current != expected and strict:
        return False, (
            f"expected_text mismatch: have {current!r}, expected {expected!r}"
        )
    runs[ri].text = new_text
    return True, f"replaced run: {current!r} -> {new_text!r}"


# ------------------------------------------------------------------
# Slide pruning: keep only selected slides, in order
# ------------------------------------------------------------------
def prune_slides(prs, selected_1indexed: list[int]) -> None:
    sldIdLst = prs.slides._sldIdLst
    sld_ids = list(sldIdLst)
    total = len(sld_ids)
    keep_idx0 = [i - 1 for i in selected_1indexed]
    for i in keep_idx0:
        if i < 0 or i >= total:
            raise ValueError(f"selected slide {i + 1} out of range (1..{total})")

    # Build new order list of element references.
    new_order = [sld_ids[i] for i in keep_idx0]

    # Remove all entries, then re-append in the new order. Keeping the same
    # element instances preserves their r:id and avoids orphaning relationships.
    for sld in list(sldIdLst):
        sldIdLst.remove(sld)
    for sld in new_order:
        sldIdLst.append(sld)


# ------------------------------------------------------------------
# Main driver
# ------------------------------------------------------------------
def run(args) -> int:
    edits_spec = json.loads(args.edits.read_text(encoding="utf-8"))
    selected = list(edits_spec.get("selected_slides") or [])
    raw_edits = list(edits_spec.get("edits") or [])
    slot_index = load_slot_index(args.detail)

    prs = Presentation(args.template)
    n_slides = len(prs.slides)
    if not selected:
        selected = list(range(1, n_slides + 1))
    for s in selected:
        if s < 1 or s > n_slides:
            print(f"ERROR: selected slide {s} out of range 1..{n_slides}", file=sys.stderr)
            return 2

    # Resolve and group edits by slide_number (in TEMPLATE order, pre-prune)
    planned: list[dict] = []
    for e in raw_edits:
        slide_num = e["slide"]
        meta = None
        if "address" in e:
            addr = dict(e["address"])
            expected = e.get("expected_text") or addr.get("expected_text")
        else:
            slot_id = e.get("slot_id")
            if slot_id is None:
                print(f"ERROR: edit {e} missing both address and slot_id", file=sys.stderr)
                return 2
            key = (slide_num, slot_id)
            if key not in slot_index:
                print(f"ERROR: slot_id {slot_id!r} not found for slide {slide_num} in detail.json",
                      file=sys.stderr)
                return 2
            meta = slot_index[key]
            addr = dict(meta["address"])
            expected = e.get("expected_text") or meta.get("expected_text")
        planned.append({
            "slide": slide_num,
            "address": addr,
            "expected": expected,
            "new_text": e["new_text"],
            "meta": meta,
            "raw": e,
        })

    if args.dry_run:
        for p in planned:
            print(f"slide {p['slide']:>3} addr={p['address']} -> {p['new_text']!r}")
        return 0

    # Apply edits BEFORE pruning so slide indices match the template
    overflow_issues: list[str] = []
    for p in planned:
        slide = prs.slides[p["slide"] - 1]
        shape_id = p["address"].get("shape_id")
        shape = find_shape(slide.shapes, shape_id)
        if shape is None:
            msg = f"slide {p['slide']}: shape_id {shape_id} not found"
            if args.strict:
                print(f"ERROR: {msg}", file=sys.stderr)
                return 3
            print(f"WARN:  {msg}", file=sys.stderr)
            continue
        ok, info = apply_edit_to_shape(
            shape, p["address"], p["new_text"], p["expected"], args.strict
        )
        prefix = "OK  " if ok else ("ERR " if args.strict else "WARN")
        print(f"{prefix} slide {p['slide']:>3} shape {shape_id}: {info}")
        if not ok and args.strict:
            return 4

        # Overflow lint (only when slot capacity metadata is available)
        if not args.no_lint and p["meta"]:
            fits, omsg = check_overflow(p["new_text"], p["meta"])
            if omsg and not fits:
                slot_id = p["meta"].get("slot_id", "?")
                line = (f"slide {p['slide']:>3} {slot_id}: 文字过长 {omsg} "
                        f"-> {p['new_text'][:24]!r}")
                overflow_issues.append(line)
                print(f"OVERFLOW {line}", file=sys.stderr)
            elif omsg and fits:  # autofit soft note
                print(f"note  slide {p['slide']:>3}: {omsg}", file=sys.stderr)

    if overflow_issues:
        print(f"\n{len(overflow_issues)} 处文字可能出框（请缩短文字到容量内，"
              f"保持字号不变以维持同级一致）：", file=sys.stderr)
        for line in overflow_issues:
            print(f"  - {line}", file=sys.stderr)
        if args.strict:
            print("\n--strict: 因出框风险拒绝保存。缩短上述文字后重试。", file=sys.stderr)
            return 5

    prune_slides(prs, selected)
    prs.save(args.output)
    note = f"\nSaved {args.output} with {len(selected)} slides and {len(planned)} edits"
    if overflow_issues:
        note += f"  (⚠️ {len(overflow_issues)} 处可能出框，见上)"
    print(note)
    return 0


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("template", type=Path)
    ap.add_argument("edits", type=Path)
    ap.add_argument("output", type=Path)
    ap.add_argument("--detail", type=Path, default=None)
    ap.add_argument("--no-lint", action="store_true",
                    help="disable the text-overflow check")
    ap.add_argument("--strict", action="store_true")
    ap.add_argument("--dry-run", action="store_true")
    args = ap.parse_args()

    if args.detail is None:
        # Try sibling detail.json next to the template
        guess = args.template.with_name("detail.json")
        if guess.exists():
            args.detail = guess
            print(f"(auto-detected detail.json: {guess})", file=sys.stderr)

    return run(args)


if __name__ == "__main__":
    sys.exit(main())
