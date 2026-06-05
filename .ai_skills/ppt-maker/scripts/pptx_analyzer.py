#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def emu(value) -> int:
    return int(value or 0)


def text_runs(shape) -> list[dict[str, object]]:
    if not getattr(shape, "has_text_frame", False):
        return []
    rows: list[dict[str, object]] = []
    for p_idx, paragraph in enumerate(shape.text_frame.paragraphs):
        for r_idx, run in enumerate(paragraph.runs):
            if not run.text:
                continue
            rows.append(
                {
                    "shape_id": shape.shape_id,
                    "paragraph": p_idx,
                    "run": r_idx,
                    "text": run.text,
                    "left": emu(shape.left),
                    "top": emu(shape.top),
                    "width": emu(shape.width),
                    "height": emu(shape.height),
                    "font_size_pt": float(run.font.size.pt) if run.font.size else None,
                    "bold": bool(run.font.bold),
                    "italic": bool(run.font.italic),
                }
            )
    return rows


def table_info(shape) -> dict[str, object] | None:
    if not getattr(shape, "has_table", False):
        return None
    table = shape.table
    return {
        "shape_id": shape.shape_id,
        "rows": len(table.rows),
        "columns": len(table.columns),
        "text_matrix": [[cell.text for cell in row.cells] for row in table.rows],
    }


def chart_info(shape) -> dict[str, object]:
    base = {
        "shape_id": shape.shape_id,
        "chart_type": "",
        "supported": True,
        "series_count": 0,
        "category_count": 0,
        "raw_tag": "",
        "error": "",
    }
    try:
        chart = shape.chart
        base["chart_type"] = str(chart.chart_type).split(" ")[0]
        base["series_count"] = len(chart.series)
        try:
            first_series = chart.series[0] if chart.series else None
            base["category_count"] = len(first_series.categories) if first_series is not None else 0
        except Exception:
            base["category_count"] = 0
    except Exception as exc:
        text = str(exc)
        match = re.search(r"unsupported plot type ([^\s]+)", text)
        base["chart_type"] = "UNSUPPORTED"
        base["supported"] = False
        base["raw_tag"] = match.group(1) if match else ""
        base["error"] = text
    return base


def image_info(shape) -> dict[str, object] | None:
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return None
    return {
        "shape_id": shape.shape_id,
        "kind": "picture",
        "left": emu(shape.left),
        "top": emu(shape.top),
        "width": emu(shape.width),
        "height": emu(shape.height),
    }


def walk_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from walk_shapes(shape.shapes)


def extract_pptx(path: Path) -> dict[str, object]:
    prs = Presentation(str(path))
    slides = []
    for slide_number, slide in enumerate(prs.slides, 1):
        texts: list[dict[str, object]] = []
        charts: list[dict[str, object]] = []
        images: list[dict[str, object]] = []
        tables: list[dict[str, object]] = []
        for shape in walk_shapes(slide.shapes):
            texts.extend(text_runs(shape))
            if getattr(shape, "has_chart", False):
                charts.append(chart_info(shape))
            table = table_info(shape)
            if table:
                tables.append(table)
            image = image_info(shape)
            if image:
                images.append(image)
        slides.append(
            {
                "slide_number": slide_number,
                "size": {"width": emu(prs.slide_width), "height": emu(prs.slide_height)},
                "texts": texts,
                "charts": charts,
                "images": images,
                "tables": tables,
            }
        )
    return {
        "schema": "ppt-maker-machine-extracted/v1",
        "source_pptx": str(path),
        "slide_count": len(prs.slides),
        "slides": slides,
    }


def inspect_pptx(path: Path) -> str:
    payload = extract_pptx(path)
    lines = [f"PPTX: {path}", f"Slides: {payload['slide_count']}"]
    for slide in payload["slides"]:
        lines.append(
            f"slide {slide['slide_number']}: texts={len(slide['texts'])} "
            f"charts={len(slide['charts'])} images={len(slide['images'])} tables={len(slide['tables'])}"
        )
    return "\n".join(lines) + "\n"


def run_cli(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Inspect or extract PPTX structure.")
    parser.add_argument("--mode", choices=("inspect", "extract"), required=True)
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--output", type=Path)
    args = parser.parse_args(argv)
    if not args.pptx.exists() or args.pptx.suffix.lower() != ".pptx":
        print("input_missing_or_not_pptx")
        return 1
    try:
        if args.mode == "inspect":
            print(inspect_pptx(args.pptx), end="")
            return 0
        payload = extract_pptx(args.pptx)
        text = json.dumps(payload, ensure_ascii=False, indent=2) + "\n"
        if args.output:
            args.output.parent.mkdir(parents=True, exist_ok=True)
            args.output.write_text(text, encoding="utf-8")
        else:
            print(text, end="")
        return 0
    except Exception as exc:
        print(f"parse_failed: {exc}")
        return 2


if __name__ == "__main__":
    raise SystemExit(run_cli())
